"""
File preprocessor for document ingestion.

Supports .docx, .pdf, .txt, .pptx, .xlsx. Extracts text, tables, and images from
all supported formats. Outputs a list of parts for downstream ingestion.
No embeddings, vector store, or LLM — ingestion is handled elsewhere.
"""

import base64
import logging
import os
import zipfile
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Final, Literal

import fitz  # PyMuPDF
from pydantic import BaseModel, Field

# -----------------------------------------------------------------------------
# Logging
# -----------------------------------------------------------------------------

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger(__name__)

# -----------------------------------------------------------------------------
# Constants & Types
# -----------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: Final[tuple[str, ...]] = (".docx", ".pdf", ".txt", ".pptx", ".xlsx")

PartType = Literal["text", "image", "table"]


def _build_table_metadata(
    source_file: str,
    table_index: int,
    row_count: int,
    *,
    page: int | None = None,
    slide: int | None = None,
    columns: list[str] | None = None,
) -> dict[str, Any]:
    """
    Build consistent table metadata for all formats.
    Always includes: source_file, table_id, table_index, row_count.
    Includes page (PDF), slide (PPTX), or neither (DOCX).
    Includes columns when available (first row as headers).
    """
    basename = os.path.basename(source_file)
    if page is not None:
        table_id = f"{basename}#page_{page}_table_{table_index}"
    elif slide is not None:
        table_id = f"{basename}#slide_{slide}_table_{table_index}"
    else:
        table_id = f"{basename}#table_{table_index}"
    meta: dict[str, Any] = {
        "source_file": source_file,
        "table_id": table_id,
        "table_index": table_index,
        "row_count": row_count,
    }
    if page is not None:
        meta["page"] = page
    if slide is not None:
        meta["slide"] = slide
    if columns is not None:
        meta["columns"] = columns
    return meta


class Part(BaseModel):
    """A single extracted part (text, table, or image) from a document."""

    type: PartType = Field(..., description="Kind of part: text, table, or image")
    content: str = Field(..., description="Raw text, table text, or base64-encoded image")
    metadata: dict[str, Any] = Field(default_factory=dict, description="Source file, page/slide, etc.")

    model_config = {"frozen": False}


# -----------------------------------------------------------------------------
# Path Resolution
# -----------------------------------------------------------------------------


class PathResolver:
    """Resolves file paths relative to cwd or project root."""

    @staticmethod
    def resolve(file_path: str, script_dir: str | None = None) -> str:
        """
        Resolve path from cwd or project root (parent of script directory).

        Args:
            file_path: Path as given (relative or absolute).
            script_dir: Directory containing the calling script; defaults to
                this module's directory.

        Returns:
            Resolved absolute path if found; otherwise the original path.
        """
        if os.path.isabs(file_path) and os.path.exists(file_path):
            return file_path
        if os.path.exists(file_path):
            return os.path.abspath(file_path)
        base = script_dir or os.path.dirname(os.path.abspath(__file__))
        project_root = os.path.dirname(base)
        candidate = os.path.join(project_root, file_path)
        if os.path.exists(candidate):
            return candidate
        return file_path


# -----------------------------------------------------------------------------
# Extractors (Strategy pattern)
# -----------------------------------------------------------------------------


class BaseExtractor(ABC):
    """Abstract base for format-specific document extractors."""

    @property
    @abstractmethod
    def extensions(self) -> tuple[str, ...]:
        """File extensions this extractor supports (e.g. ('.pdf',))."""
        ...

    @abstractmethod
    def extract(self, file_path: str) -> list[Part]:
        """
        Extract text and images from the file into a list of parts.

        Args:
            file_path: Absolute path to the file.

        Returns:
            List of Part dicts (type, content, metadata).
        """
        ...

    def supports(self, extension: str) -> bool:
        """Return True if this extractor handles the given extension."""
        return extension.lower() in self.extensions


class PdfExtractor(BaseExtractor):
    """Extract text, tables, and images from PDF files (PyMuPDF + pdfplumber)."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return (".pdf",)

    def extract(self, file_path: str) -> list[Part]:
        parts: list[Part] = []
        tables_by_page = self._extract_tables_by_page(file_path)
        doc = fitz.open(file_path)
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                self._extract_page_text(page, page_num, file_path, parts)
                for table_idx, (content, row_count, columns) in enumerate(
                    tables_by_page.get(page_num, [])
                ):
                    metadata = _build_table_metadata(
                        file_path,
                        table_idx,
                        row_count,
                        page=page_num,
                        columns=columns,
                    )
                    parts.append(Part(type="table", content=content, metadata=metadata))
                self._extract_page_images(doc, page, page_num, file_path, parts)
        finally:
            doc.close()
        return parts

    def _extract_page_text(
        self,
        page: fitz.Page,
        page_num: int,
        file_path: str,
        parts: list[Part],
    ) -> None:
        """Extract text from a single page using PyMuPDF."""
        text = page.get_text().strip()
        if text:
            parts.append(
                Part(
                    type="text",
                    content=text,
                    metadata={"source_file": file_path, "page": page_num},
                )
            )

    def _extract_page_images(
        self,
        doc: fitz.Document,
        page: fitz.Page,
        page_num: int,
        file_path: str,
        parts: list[Part],
    ) -> None:
        """Extract images from a single page using PyMuPDF."""
        for img_index, img in enumerate(page.get_images(full=True)):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                b64 = base64.b64encode(image_bytes).decode()
                image_id = f"page_{page_num}_img_{img_index}"
                parts.append(
                    Part(
                        type="image",
                        content=b64,
                        metadata={
                            "source_file": file_path,
                            "page": page_num,
                            "image_id": image_id,
                        },
                    )
                )
            except Exception as e:
                logger.warning(
                    "Skip image %s on page %s: %s",
                    img_index,
                    page_num,
                    e,
                    exc_info=False,
                )

    def _extract_tables_by_page(
        self, file_path: str
    ) -> dict[int, list[tuple[str, int, list[str] | None]]]:
        """
        Extract tables from PDF using pdfplumber.
        Returns page_num -> list of (content, row_count, columns).
        columns is first row as list of strings, or None.
        """
        out: dict[int, list[tuple[str, int, list[str] | None]]] = {}
        try:
            import pdfplumber
        except ImportError:
            logger.debug("pdfplumber not installed; skipping PDF table extraction")
            return out
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num in range(len(pdf.pages)):
                    page = pdf.pages[page_num]
                    tables = page.extract_tables()
                    if not tables:
                        continue
                    entries: list[tuple[str, int, list[str] | None]] = []
                    for table in tables:
                        if not table:
                            continue
                        row_count = len(table)
                        cells_per_row = [" | ".join((str(cell or "").strip().replace("\n", " ") for cell in row)) for row in table]
                        if not cells_per_row:
                            continue
                        content = "Table:\n" + "\n".join(cells_per_row)
                        columns = [str(cell or "").strip() for cell in table[0]] if table else None
                        entries.append((content, row_count, columns))
                    if entries:
                        out[page_num] = entries
        except Exception as e:
            logger.warning("PDF table extraction failed: %s", e, exc_info=False)
        return out


class PptxExtractor(BaseExtractor):
    """Extract text and images from PowerPoint (.pptx) files."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return (".pptx",)

    def extract(self, file_path: str) -> list[Part]:
        from pptx import Presentation

        parts: list[Part] = []
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides):
            slide_images = self._collect_slide_content(slide, slide_idx, file_path, parts)
            for b64, image_id in slide_images:
                parts.append(
                    Part(
                        type="image",
                        content=b64,
                        metadata={
                            "source_file": file_path,
                            "slide": slide_idx,
                            "image_id": image_id,
                        },
                    )
                )
        return parts

    def _collect_slide_content(
        self,
        slide: Any,
        slide_idx: int,
        file_path: str,
        parts: list[Part],
    ) -> list[tuple[str, str]]:
        """Collect text and table parts, append image (b64, image_id) list. Returns slide_images."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide_texts: list[str] = []
        slide_images: list[tuple[str, str]] = []
        title_shape = slide.shapes.title
        if title_shape and title_shape.text.strip():
            slide_texts.append("Slide Title: " + title_shape.text.strip())
        img_count = 0
        table_index = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape is not title_shape:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                tbl = shape.table
                rows = [
                    " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                    for row in tbl.rows
                ]
                if rows:
                    row_count = len(rows)
                    columns = [cell.text.strip() for cell in tbl.rows[0].cells] if tbl.rows else None
                    metadata = _build_table_metadata(
                        file_path,
                        table_index,
                        row_count,
                        slide=slide_idx,
                        columns=columns,
                    )
                    content = "Table:\n" + "\n".join(rows)
                    parts.append(Part(type="table", content=content, metadata=metadata))
                    table_index += 1
            if (
                getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                and hasattr(shape, "image")
            ):
                b64 = base64.b64encode(shape.image.blob).decode()
                image_id = f"slide_{slide_idx}_img_{img_count}"
                img_count += 1
                slide_images.append((b64, image_id))
        if slide_texts:
            parts.append(
                Part(
                    type="text",
                    content="\n".join(slide_texts),
                    metadata={"source_file": file_path, "slide": slide_idx},
                )
            )
        return slide_images


class DocxExtractor(BaseExtractor):
    """Extract text, tables, and images from Word (.docx) files."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return (".docx",)

    def extract(self, file_path: str) -> list[Part]:
        try:
            from docx import Document as DocxDocument
        except ImportError as e:
            logger.error("python-docx required for .docx. pip install python-docx")
            raise e

        parts: list[Part] = []
        doc = DocxDocument(file_path)
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                parts.append(
                    Part(
                        type="text",
                        content=text,
                        metadata={"source_file": file_path, "paragraph": i},
                    )
                )
        for i, table in enumerate(doc.tables):
            rows = [
                " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                for row in table.rows
            ]
            if rows:
                row_count = len(rows)
                columns = [cell.text.strip() for cell in table.rows[0].cells] if table.rows else None
                metadata = _build_table_metadata(
                    file_path,
                    i,
                    row_count,
                    columns=columns,
                )
                parts.append(
                    Part(
                        type="table",
                        content="Table:\n" + "\n".join(rows),
                        metadata=metadata,
                    )
                )
        self._extract_images(file_path, parts)
        return parts

    def _extract_images(self, file_path: str, parts: list[Part]) -> None:
        """Extract embedded images from DOCX (word/media in the zip)."""
        try:
            with zipfile.ZipFile(file_path, "r") as zf:
                idx = 0
                for name in zf.namelist():
                    if not name.startswith("word/media/") or name == "word/media/":
                        continue
                    try:
                        data = zf.read(name)
                        b64 = base64.b64encode(data).decode()
                        image_id = f"docx_img_{idx}"
                        idx += 1
                        parts.append(
                            Part(
                                type="image",
                                content=b64,
                                metadata={
                                    "source_file": file_path,
                                    "image_id": image_id,
                                    "media_path": name,
                                },
                            )
                        )
                    except Exception as e:
                        logger.warning("Skip DOCX image %s: %s", name, e)
        except Exception as e:
            logger.warning("DOCX image extraction failed: %s", e, exc_info=False)


class PlainTextExtractor(BaseExtractor):
    """Extract content from plain text (.txt) files."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return (".txt",)

    def extract(self, file_path: str) -> list[Part]:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        if not content.strip():
            return []
        return [
            Part(
                type="text",
                content=content,
                metadata={"source_file": file_path},
            )
        ]


class XlsxExtractor(BaseExtractor):
    """Extract tables and data from Excel (.xlsx) files."""

    @property
    def extensions(self) -> tuple[str, ...]:
        return (".xlsx",)

    def extract(self, file_path: str) -> list[Part]:
        try:
            import pandas as pd
            import openpyxl
        except ImportError as e:
            logger.error("pandas and openpyxl required for .xlsx. pip install pandas openpyxl")
            raise e

        parts: list[Part] = []
        try:
            # Get sheet names
            xl = pd.ExcelFile(file_path, engine='openpyxl')
            sheet_names = xl.sheet_names

            for sheet_idx, sheet_name in enumerate(sheet_names):
                try:
                    df = pd.read_excel(xl, sheet_name=sheet_name)
                    if df.empty:
                        continue

                    # Get column names
                    columns = [str(col).strip() for col in df.columns.tolist()]
                    row_count = len(df)

                    # Convert DataFrame to table format
                    rows = []
                    # Add header row
                    rows.append(" | ".join(columns))
                    # Add data rows
                    for _, row in df.iterrows():
                        row_values = [str(val).strip().replace("\n", " ") if pd.notna(val) else "" for val in row]
                        rows.append(" | ".join(row_values))

                    content = "Table:\n" + "\n".join(rows)

                    # Build metadata
                    basename = os.path.basename(file_path)
                    table_id = f"{basename}#sheet_{sheet_name}"
                    metadata = {
                        "source_file": file_path,
                        "table_id": table_id,
                        "table_index": sheet_idx,
                        "row_count": row_count,
                        "sheet_name": sheet_name,
                        "sheet_index": sheet_idx,
                        "columns": columns,
                        "column_count": len(columns),
                    }

                    parts.append(Part(type="table", content=content, metadata=metadata))
                    logger.info("Extracted sheet '%s' with %d rows", sheet_name, row_count)

                except Exception as e:
                    logger.warning("Failed to extract sheet '%s': %s", sheet_name, e)
                    continue

        except Exception as e:
            logger.error("Excel extraction failed: %s", e, exc_info=False)

        return parts


# -----------------------------------------------------------------------------
# File Preprocessor (Facade)
# -----------------------------------------------------------------------------


class FilePreprocessor:
    """
    Main entry point for file preprocessing.

    Dispatches to format-specific extractors and returns a list of parts
    (text and images) for downstream ingestion.
    """

    def __init__(self, path_resolver: PathResolver | None = None) -> None:
        """
        Initialize the preprocessor with default extractors.

        Args:
            path_resolver: Optional custom path resolver; defaults to PathResolver().
        """
        self._path_resolver = path_resolver or PathResolver()
        self._extractors: dict[str, BaseExtractor] = {}
        self._register_default_extractors()

    def _register_default_extractors(self) -> None:
        """Register built-in extractors for all supported formats."""
        for extractor in (
            PdfExtractor(),
            PptxExtractor(),
            DocxExtractor(),
            PlainTextExtractor(),
            XlsxExtractor(),
        ):
            for ext in extractor.extensions:
                self._extractors[ext] = extractor

    def preprocess(self, file_path: str, *, resolve: bool = True) -> list[Part]:
        """
        Preprocess a file and return a list of parts.

        Args:
            file_path: Path to the file (relative or absolute).
            resolve: If True, resolve path from cwd or project root.

        Returns:
            List of Part models. Each has type ("text", "table", or "image"),
            content (raw text, table text, or base64 image), and metadata.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file extension is not supported.
        """
        path = (
            self._path_resolver.resolve(file_path)
            if resolve
            else os.path.abspath(file_path)
        )
        if not os.path.exists(path):
            raise FileNotFoundError(f"File not found: {path}")

        ext = Path(path).suffix.lower()
        if ext not in self._extractors:
            raise ValueError(
                f"Unsupported extension '{ext}'. "
                f"Supported: {', '.join(sorted(self._extractors))}"
            )

        logger.info("Preprocessing %s (%s)", path, ext)
        extractor = self._extractors[ext]
        parts = extractor.extract(path)
        logger.info("Extracted %d parts from %s", len(parts), path)
        return parts

    def is_supported(self, file_path: str) -> bool:
        """Return True if the file extension is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in self._extractors

    def get_supported_extensions(self) -> tuple[str, ...]:
        """Return all supported file extensions."""
        return tuple(sorted(self._extractors.keys()))


# -----------------------------------------------------------------------------
# Module-level API (backward compatibility)
# -----------------------------------------------------------------------------

_default_preprocessor = FilePreprocessor()


def preprocess(file_path: str, resolve: bool = True) -> list[Part]:
    """
    Preprocess a file and return a list of parts.

    Convenience wrapper around the default FilePreprocessor. Supported
    formats: .docx, .pdf, .txt, .md, .pptx. Text, tables, and images are
    extracted from PDF, DOCX, and PPTX; .txt and .md return a single text part.
    """
    return _default_preprocessor.preprocess(file_path, resolve=resolve)


def is_supported(file_path: str) -> bool:
    """Return True if the file extension is supported."""
    return _default_preprocessor.is_supported(file_path)


# -----------------------------------------------------------------------------
# CLI
# -----------------------------------------------------------------------------


def _main() -> int:
    """Run the preprocessor from the command line. Returns exit code."""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ingester_agent.py <file_path>")
        print(f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}")
        return 1

    raw_path = sys.argv[1]
    try:
        path = PathResolver.resolve(raw_path)
        if not os.path.exists(path):
            logger.error("File not found: %s", raw_path)
            return 1
        parts = _default_preprocessor.preprocess(path, resolve=False)
        for i, part in enumerate(parts):
            kind = part.type
            content = part.content
            preview = (
                (content[:80] + "…") if kind == "text" and len(content) > 80 else content[:80]
            )
            print(f"  [{i}] {kind}: {preview}  metadata={part.metadata}")
        return 0
    except (FileNotFoundError, ValueError) as e:
        logger.error("%s", e)
        return 1
    except Exception:
        logger.exception("Preprocessing failed")
        return 1


if __name__ == "__main__":
    raise SystemExit(_main())
